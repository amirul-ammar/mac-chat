"""Text extraction from common file types. All parsing happens locally."""
from __future__ import annotations

from pathlib import Path

PLAIN_EXTS = {
    ".txt", ".md", ".markdown", ".rst", ".log", ".csv", ".tsv", ".json",
    ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".xml", ".html", ".htm",
    ".css", ".scss", ".sass", ".less", ".tex", ".srt", ".vtt", ".env.example",
    ".py", ".js", ".jsx", ".ts", ".tsx", ".java", ".c", ".h", ".cpp", ".hpp",
    ".cs", ".go", ".rs", ".rb", ".php", ".swift", ".kt", ".kts", ".scala",
    ".sh", ".bash", ".zsh", ".fish", ".sql", ".r", ".m", ".mm", ".lua",
    ".pl", ".vue", ".svelte", ".dart", ".ipynb", ".gitignore", ".plist",
}

RICH_EXTS = {".pdf", ".docx", ".xlsx", ".pptx", ".rtf"}

EXTRACTABLE = PLAIN_EXTS | RICH_EXTS


def can_extract(ext: str) -> bool:
    return ext.lower() in EXTRACTABLE


def _plain(path: Path, limit: int) -> str:
    with path.open("r", encoding="utf-8", errors="replace") as fh:
        return fh.read(limit)


def _pdf(path: Path, limit: int) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts: list[str] = []
    total = 0
    for page in reader.pages[:40]:
        try:
            t = page.extract_text() or ""
        except Exception:
            continue
        parts.append(t)
        total += len(t)
        if total >= limit:
            break
    return "\n".join(parts)[:limit]


def _docx(path: Path, limit: int) -> str:
    import docx

    doc = docx.Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(c.text.strip() for c in row.cells))
    return "\n".join(parts)[:limit]


def _xlsx(path: Path, limit: int) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    total = 0
    for ws in wb.worksheets:
        parts.append(f"# sheet: {ws.title}")
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i > 300 or total >= limit:
                break
            line = " | ".join("" if v is None else str(v) for v in row)
            if line.strip(" |"):
                parts.append(line)
                total += len(line)
        if total >= limit:
            break
    wb.close()
    return "\n".join(parts)[:limit]


def _pptx(path: Path, limit: int) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"# slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
    return "\n".join(parts)[:limit]


def _rtf(path: Path, limit: int) -> str:
    import re

    raw = path.read_text(encoding="utf-8", errors="replace")[: limit * 3]
    raw = re.sub(r"\\'[0-9a-fA-F]{2}", "", raw)
    raw = re.sub(r"\{\\\*?[^{}]*\}", " ", raw)
    raw = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", raw)
    raw = raw.replace("{", " ").replace("}", " ")
    return re.sub(r"\s+", " ", raw).strip()[:limit]


_HANDLERS = {
    ".pdf": _pdf,
    ".docx": _docx,
    ".xlsx": _xlsx,
    ".pptx": _pptx,
    ".rtf": _rtf,
}


def extract(path: Path, limit: int = 200_000) -> str | None:
    """Return extracted text, or None if the file yielded nothing usable."""
    ext = path.suffix.lower()
    try:
        handler = _HANDLERS.get(ext)
        text = handler(path, limit) if handler else (_plain(path, limit) if ext in PLAIN_EXTS else None)
    except Exception:
        return None
    if not text:
        return None
    text = text.strip()
    return text or None
